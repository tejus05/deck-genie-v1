import streamlit as st
import io
from typing import Dict, Any, List, Optional, Set
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from image_fetcher import fetch_image_for_slide, get_slide_icon
from utils import FONTS, COLORS, FONT_SIZES, MARGINS, CONTENT_AREA, IMAGE_AREA

def create_market_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the market opportunity slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "") if presentation_context else ""
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    
    # Replace any placeholder text
    title_p.text = content.get("title", "Market Opportunity").replace("[Product Name]", product_name)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Market size - highlighted in top-left section
    market_size_left = Inches(0.5)
    market_size_top = Inches(1.8)
    market_size_width = Inches(3.5)
    market_size_height = Inches(1.5)
    
    # Create the market size box with border - using "white" instead of "light_blue"
    market_size_box = slide.shapes.add_shape(1, market_size_left, market_size_top, market_size_width, market_size_height)
    market_size_box.fill.solid()
    market_size_box.fill.fore_color.rgb = RGBColor.from_string(COLORS["white"])  # Using white instead of light_blue
    market_size_box.line.color.rgb = RGBColor.from_string(COLORS["blue"])
    market_size_box.line.width = Pt(1.5)
    
    # Add market size text
    market_size_text_box = slide.shapes.add_textbox(market_size_left + Inches(0.1), market_size_top + Inches(0.1), 
                                                   market_size_width - Inches(0.2), market_size_height - Inches(0.2))
    market_size_tf = market_size_text_box.text_frame
    market_size_tf.word_wrap = True
    
    # Add header paragraph
    header_p = market_size_tf.paragraphs[0]
    header_p.text = "Market Size"
    header_p.alignment = PP_ALIGN.CENTER
    
    for run in header_p.runs:
        run.font.name = FONTS["body"]
        run.font.size = Pt(14)
        run.font.bold = True
        
    # Add value paragraph
    value_p = market_size_tf.add_paragraph()
    value_p.text = content.get("market_size", "$25B by 2025")
    value_p.alignment = PP_ALIGN.CENTER
    value_p.space_before = Pt(10)
    
    for run in value_p.runs:
        run.font.name = FONTS["body"]
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["blue"])
    
    # Growth rate - highlighted in top-right section
    growth_rate_left = Inches(4.5)
    growth_rate_top = market_size_top
    growth_rate_width = market_size_width
    growth_rate_height = market_size_height
    
    # Create the growth rate box with border - using "white" instead of "light_blue"
    growth_rate_box = slide.shapes.add_shape(1, growth_rate_left, growth_rate_top, growth_rate_width, growth_rate_height)
    growth_rate_box.fill.solid()
    growth_rate_box.fill.fore_color.rgb = RGBColor.from_string(COLORS["white"])  # Using white instead of light_blue
    growth_rate_box.line.color.rgb = RGBColor.from_string(COLORS["blue"])
    growth_rate_box.line.width = Pt(1.5)
    
    # Add growth rate text
    growth_rate_text_box = slide.shapes.add_textbox(growth_rate_left + Inches(0.1), growth_rate_top + Inches(0.1), 
                                                   growth_rate_width - Inches(0.2), growth_rate_height - Inches(0.2))
    growth_rate_tf = growth_rate_text_box.text_frame
    growth_rate_tf.word_wrap = True
    
    # Add header paragraph
    header_p = growth_rate_tf.paragraphs[0]
    header_p.text = "Growth Rate"
    header_p.alignment = PP_ALIGN.CENTER
    
    for run in header_p.runs:
        run.font.name = FONTS["body"]
        run.font.size = Pt(14)
        run.font.bold = True
        
    # Add value paragraph
    value_p = growth_rate_tf.add_paragraph()
    value_p.text = content.get("growth_rate", "15% CAGR")
    value_p.alignment = PP_ALIGN.CENTER
    value_p.space_before = Pt(10)
    
    for run in value_p.runs:
        run.font.name = FONTS["body"]
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["blue"])
    
    # Description text - positioned in bottom half
    desc_left = Inches(0.5)
    desc_top = Inches(3.6)
    desc_width = Inches(7.5)
    desc_height = Inches(3.0)
    
    desc_box = slide.shapes.add_textbox(desc_left, desc_top, desc_width, desc_height)
    desc_tf = desc_box.text_frame
    desc_tf.word_wrap = True
    desc_tf.auto_size = False
    
    desc_p = desc_tf.paragraphs[0]
    desc_p.text = content.get("description", "The market is experiencing significant growth as organizations increasingly recognize the need for innovative solutions in this space.")
    desc_p.alignment = PP_ALIGN.LEFT
    
    # Apply description formatting with proper spacing
    for run in desc_p.runs:
        run.font.name = FONTS["body"]
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Add image with consistent positioning
    if used_images is None:
        used_images = set()
    
    try:
        from ppt_generator import fetch_image_with_cache
        image_data = fetch_image_with_cache("market", presentation_context, used_images, use_placeholders=False)
    except ImportError:
        # Fallback to direct fetch if fetch_image_with_cache is not available
        image_data = fetch_image_for_slide("market", presentation_context, use_placeholders=False)
    
    if image_data:
        # Position image on right side
        pic = slide.shapes.add_picture(
            image_data,
            Inches(8.3),
            Inches(3.0),
            Inches(4.5),
            Inches(4.0)
        )
    else:
        # Add fallback icon
        icon_left = Inches(9.5)
        icon_top = Inches(4.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("market")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        
        # Apply icon formatting
        for run in icon_p.runs:
            run.font.name = FONTS["body"]
            run.font.size = Pt(64)
            run.font.color.rgb = RGBColor.from_string(COLORS["black"])

def create_roadmap_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """
    Create a roadmap slide with horizontal timeline layout.
    This design avoids text overlapping by using clear horizontal separation of phases.
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "") if presentation_context else ""
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    
    # Handle different title fields
    title_text = content.get("title", f"{product_name} Roadmap")
    title_p.text = title_text.replace("[Product Name]", product_name)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Get and handle phases from content
    phases = content.get("phases", [])
    if not phases:
        # Create default phases if none exist
        phases = [
            {"name": "Phase 1: Foundation", "items": ["Initial launch", "Core features", "First customers"]},
            {"name": "Phase 2: Expansion", "items": ["Advanced analytics", "Integration APIs", "Expanded support"]},
            {"name": "Phase 3: Evolution", "items": ["Enterprise features", "Global expansion", "Industry partnerships"]}
        ]
    
    # Limit to 3 phases max to prevent crowding
    phases = phases[:3]
    phase_count = len(phases)
    
    # First add phase headers with dates (if any) above the timeline
    phase_headers_top = Inches(1.7)
    phase_headers_height = Inches(0.8)
    timeline_width = Inches(10.33)  # Leave margin on both sides
    
    for i, phase in enumerate(phases):
        # Calculate position for this phase
        phase_width = timeline_width / phase_count
        phase_left = Inches(1.5) + (i * phase_width)
        
        # Add phase header textbox
        header_box = slide.shapes.add_textbox(phase_left, phase_headers_top, phase_width, phase_headers_height)
        header_tf = header_box.text_frame
        header_tf.word_wrap = True
        header_tf.auto_size = False
        header_p = header_tf.paragraphs[0]
        
        # Format phase name with appropriate font
        phase_name = phase.get("name", f"Phase {i+1}")
        header_p.text = phase_name
        header_p.alignment = PP_ALIGN.CENTER
        header_p.font.bold = True
        header_p.font.name = FONTS["body"]
        header_p.font.size = Pt(14)
        
    # Now add the timeline bar
    timeline_top = Inches(2.5)
    timeline_left = Inches(1.5)
    timeline_height = Inches(0.12)
    
    timeline_bar = slide.shapes.add_shape(1, timeline_left, timeline_top, timeline_width, timeline_height)
    timeline_bar.fill.solid()
    timeline_bar.fill.fore_color.rgb = RGBColor.from_string(COLORS["blue"])
    timeline_bar.line.fill.background()
    
    # Add phase markers on the timeline
    marker_size = Inches(0.25)
    
    for i in range(phase_count):
        # Calculate position for this marker (centered in each phase section)
        phase_width = timeline_width / phase_count
        marker_left = timeline_left + (i * phase_width) + (phase_width / 2) - (marker_size / 2)
        marker_top = timeline_top - (marker_size / 2) + (timeline_height / 2)
        
        # Create square marker
        marker = slide.shapes.add_shape(1, marker_left, marker_top, marker_size, marker_size)
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor.from_string(COLORS["blue"])
        marker.line.fill.background()
    
    # Now add bullet points for each phase below the timeline
    bullets_top = Inches(2.9)  # Start below timeline
    bullets_height = Inches(3.0)  # Adequate room for all bullets
    
    for i, phase in enumerate(phases):
        # Calculate position for this phase's bullets
        phase_width = timeline_width / phase_count
        bullets_left = timeline_left + (i * phase_width) + Inches(0.2)  # Small indent
        bullets_width = phase_width - Inches(0.4)  # Leave some margin between columns
        
        # Add bullet points text box
        bullets_box = slide.shapes.add_textbox(bullets_left, bullets_top, bullets_width, bullets_height)
        bullets_tf = bullets_box.text_frame
        bullets_tf.word_wrap = True
        bullets_tf.auto_size = False
        
        # Get items for this phase
        items = phase.get("items", [])
        items = items[:4]  # Limit to 4 items max per phase
        
        # Add items as bullet points
        for j, item in enumerate(items):
            if j == 0:
                p = bullets_tf.paragraphs[0]
            else:
                p = bullets_tf.add_paragraph()
            
            # Add bullet point text
            p.text = f"• {item}"
            p.alignment = PP_ALIGN.LEFT
            
            # Add space between bullets
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            
            # Format text
            for run in p.runs:
                run.font.name = FONTS["body"]
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor.from_string(COLORS["black"])

def create_team_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the team overview slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    company_name = presentation_context.get("metadata", {}).get("company_name", "") if presentation_context else ""
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    
    # Use proper title from content or generate one
    title_p.text = content.get("title", f"{company_name} Leadership Team")
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Get team members from content
    team_members = content.get("team_members", [
        {"name": "Alex Johnson", "role": "Chief Executive Officer"},
        {"name": "Sam Washington", "role": "Chief Technology Officer"},
        {"name": "Jordan Smith", "role": "VP of Product"},
        {"name": "Taylor Rivera", "role": "VP of Customer Success"}
    ])
    
    # Limit to 4 team members
    team_members = team_members[:4]
    
    # Calculate layout for team members
    left_margin = Inches(0.5)
    top_margin = Inches(1.8)
    
    # Position team members in a 2x2 grid
    person_width = Inches(6)
    person_height = Inches(2)
    
    for i, person in enumerate(team_members):
        # Calculate position (2x2 grid)
        col = i % 2
        row = i // 2
        
        person_left = left_margin + (col * person_width)
        person_top = top_margin + (row * person_height)
        
        # Create box for the person
        person_box = slide.shapes.add_textbox(person_left, person_top, person_width, person_height)
        person_tf = person_box.text_frame
        
        # Add name
        name_p = person_tf.paragraphs[0]
        name_p.text = person.get("name", f"Team Member {i+1}")
        name_p.alignment = PP_ALIGN.LEFT
        name_p.space_after = Pt(0)
        
        for run in name_p.runs:
            run.font.name = FONTS["body"]
            run.font.size = Pt(20)
            run.font.bold = True
        
        # Add role
        role_p = person_tf.add_paragraph()
        role_p.text = person.get("role", "Role")
        role_p.alignment = PP_ALIGN.LEFT
        role_p.space_before = Pt(4)
        
        for run in role_p.runs:
            run.font.name = FONTS["body"]
            run.font.size = Pt(16)
            run.font.italic = True
    
    # Add tagline at the bottom
    tagline = content.get("tagline", "Building innovative solutions for tomorrow's challenges")
    
    tagline_left = Inches(0.5)
    tagline_top = Inches(6.0)
    tagline_width = Inches(12.33)
    tagline_height = Inches(0.8)
    
    tagline_box = slide.shapes.add_textbox(tagline_left, tagline_top, tagline_width, tagline_height)
    tagline_tf = tagline_box.text_frame
    tagline_tf.word_wrap = True
    tagline_tf.auto_size = False
    tagline_p = tagline_tf.paragraphs[0]
    tagline_p.text = tagline
    tagline_p.alignment = PP_ALIGN.CENTER
    
    # Apply tagline formatting
    for run in tagline_p.runs:
        run.font.name = FONTS["body"]
        run.font.size = Pt(16)
        run.font.italic = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])