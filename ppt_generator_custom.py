import io
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from image_fetcher import get_slide_icon
from utils import FONTS, COLORS, FONT_SIZES, MARGINS, CONTENT_AREA, IMAGE_AREA, match_icon_to_feature
from image_manager import ImageManager

def create_custom_presentation(content: Dict[str, Any], filename: str, slide_order: List[str], image_manager: ImageManager) -> io.BytesIO:
    """
    Create a customized PowerPoint presentation from modified content.
    
    Args:
        content: Dictionary containing modified content for slides
        filename: Name to give the file
        slide_order: Custom order of slides
        image_manager: Manager for custom images
    Returns:
        BytesIO object containing the presentation
    """
    prs = Presentation()
    
    # Set slide size to widescreen (16:9) - 13.33 x 7.5 inches
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Create slides in custom order
    slide_creators = {
        'title_slide': lambda: create_custom_title_slide(prs, content['title_slide'], image_manager),
        'problem_slide': lambda: create_custom_problem_slide(prs, content['problem_slide'], content, image_manager),
        'solution_slide': lambda: create_custom_solution_slide(prs, content['solution_slide'], content, image_manager),
        'features_slide': lambda: create_custom_features_slide(prs, content['features_slide'], image_manager),
        'advantage_slide': lambda: create_custom_advantage_slide(prs, content['advantage_slide'], content, image_manager),
        'audience_slide': lambda: create_custom_audience_slide(prs, content['audience_slide'], content, image_manager),
        'cta_slide': lambda: create_custom_cta_slide(prs, content['cta_slide'], image_manager),
        'market_slide': lambda: create_custom_market_slide(prs, content['market_slide'], content, image_manager),
        'roadmap_slide': lambda: create_custom_roadmap_slide(prs, content['roadmap_slide'], content, image_manager),
        'team_slide': lambda: create_custom_team_slide(prs, content['team_slide'], content, image_manager)
    }
    
    # Create slides in the specified order
    for slide_key in slide_order:
        if slide_key in content and slide_key in slide_creators:
            slide_creators[slide_key]()
    
    # Save to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output

def apply_text_formatting(text_frame, font_name=FONTS["body"], size=FONT_SIZES["body"], 
                         bold=False, color=COLORS["black"], alignment=PP_ALIGN.LEFT):
    """Apply consistent text formatting to a text frame."""
    text_frame.paragraphs[0].alignment = alignment
    
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor.from_string(color)

def create_custom_title_slide(prs: Presentation, content: Dict[str, str], image_manager: ImageManager):
    """Create the customized title slide."""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background to white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(COLORS["white"])
    
    # Title
    title = slide.shapes.title
    title.text = content["title"]
    apply_text_formatting(
        title.text_frame, 
        font_name=FONTS["title"], 
        size=FONT_SIZES["title_slide"], 
        bold=True, 
        alignment=PP_ALIGN.CENTER
    )
    
    # Company name (subtitle)
    subtitle = slide.placeholders[1]
    subtitle.text = content["subtitle"]
    apply_text_formatting(
        subtitle.text_frame, 
        font_name=FONTS["body"], 
        size=FONT_SIZES["subtitle"], 
        bold=False, 
        alignment=PP_ALIGN.CENTER
    )
    
    # Add custom image if available (cached or uploaded)
    custom_image = image_manager.get_image_for_slide('title_slide')
    if custom_image:
        # Add custom image as background or decorative element
        # Position it in the lower right corner - adjusted for 16:9 format
        pic = slide.shapes.add_picture(
            custom_image,
            Inches(9.83),  # Further right for 16:9 format
            Inches(5),
            Inches(2.5),
            Inches(2)
        )

def create_custom_problem_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any], image_manager: ImageManager):
    """Create the customized problem statement slide."""
    # Use blank layout to avoid unwanted placeholder elements
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title manually
    title_left = Inches(1.0)
    title_top = Inches(0.5)
    title_width = Inches(11.33)
    title_height = Inches(1.0)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = content["title"]
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting and allow natural wrapping
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(FONT_SIZES["title"])
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Ensure title wraps naturally if it's too long
    title_tf.word_wrap = True
    
    # Create custom text box for bullet points with better spacing
    left = Inches(CONTENT_AREA["left"])
    top = Inches(CONTENT_AREA["top"])
    width = Inches(CONTENT_AREA["width"])  # Use defined width to prevent overflow into image area
    height = Inches(CONTENT_AREA["height"])
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_right = Inches(0.1)  # Add right margin to prevent text overflow
    
    # Clear any existing text and add bullet points
    if tf.paragraphs:
        p = tf.paragraphs[0]
        p.text = ""
    else:
        p = tf.add_paragraph()
        
    # Add bullet points with proper formatting and bullet characters
    for i, bullet in enumerate(content["bullets"]):
        # Truncate bullet points if they're too long (prevents text overflow)
        bullet_text = bullet
        if len(bullet) > 120:  # Reasonable character limit for bullet points
            bullet_text = bullet[:117] + "..."
            
        if i == 0:
            p.text = f"• {bullet_text}"  # Add bullet character
        else:
            p = tf.add_paragraph()
            p.text = f"• {bullet_text}"  # Add bullet character
        
        # Apply bullet formatting
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        
        # Apply font formatting
        p.font.name = FONTS["body"]
        p.font.size = Pt(FONT_SIZES["body"])
        
        # Apply spacing for better readability
        p.space_before = Pt(6)  # Add space before each bullet
        p.space_after = Pt(6)   # Add space after each bullet
        
    apply_text_formatting(tf)
    
    # Add custom image directly to slide without using placeholder
    image_data = image_manager.get_image_for_slide('problem_slide')
    
    if image_data:
        # Position the image on the right side of the 16:9 slide with proper dimensions
        pic = slide.shapes.add_picture(
            image_data,
            Inches(IMAGE_AREA["left"]),     # Positioned further right for 16:9 format
            Inches(IMAGE_AREA["top"]),
            Inches(IMAGE_AREA["width"]),    # Adjusted width for 16:9
            Inches(IMAGE_AREA["height"])
        )
    else:
        # Add fallback icon as a separate shape
        icon_left = Inches(IMAGE_AREA["left"] + IMAGE_AREA["width"]/2 - 0.5)
        icon_top = Inches(IMAGE_AREA["top"] + IMAGE_AREA["height"]/2 - 0.5)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(1), Inches(1))
        icon_tf = icon_box.text_frame
        icon_p = icon_tf.paragraphs[0]
        icon_info = get_slide_icon("problem")
        icon_p.text = icon_info["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=72, alignment=PP_ALIGN.CENTER)

def create_custom_solution_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any], image_manager: ImageManager):
    """Create the customized solution overview slide."""
    # Use blank layout to avoid unwanted placeholder elements
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title manually
    title_left = Inches(1.0)
    title_top = Inches(0.5)
    title_width = Inches(11.33)
    title_height = Inches(1.0)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = content["title"]
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting and allow natural wrapping
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(FONT_SIZES["title"])
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Ensure title wraps naturally if it's too long
    title_tf.word_wrap = True
    
    # Create custom text box for paragraph content with better spacing
    left = Inches(CONTENT_AREA["left"])
    top = Inches(CONTENT_AREA["top"])
    width = Inches(CONTENT_AREA["width"])  # Use defined width to prevent overflow into image area
    height = Inches(CONTENT_AREA["height"])
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_right = Inches(0.1)  # Add right margin to prevent text overflow
    
    p = tf.paragraphs[0]
    # Truncate paragraph if it's too long to prevent text overflow
    # Handle different field names from solution slide content generator
    paragraph_text = ""
    if "paragraph" in content:
        paragraph_text = content["paragraph"]
    elif "description" in content:
        paragraph_text = content["description"]
    elif "value_proposition" in content:
        paragraph_text = content["value_proposition"]
    
    if len(paragraph_text) > 550:  # Reasonable limit for paragraphs
        paragraph_text = paragraph_text[:547] + "..."
    
    p.text = paragraph_text
    p.alignment = PP_ALIGN.LEFT
    
    # Format paragraph with proper spacing
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    p.line_spacing = 1.2  # Better line spacing
    
    apply_text_formatting(tf)
    
    # Add custom image directly to slide without using placeholder
    image_data = image_manager.get_image_for_slide('solution_slide')
    
    if image_data:
        # Position the image on the right side of the 16:9 slide with proper dimensions
        pic = slide.shapes.add_picture(
            image_data,
            Inches(IMAGE_AREA["left"]),     # Positioned further right for 16:9 format
            Inches(IMAGE_AREA["top"]),
            Inches(IMAGE_AREA["width"]),    # Adjusted width for 16:9
            Inches(IMAGE_AREA["height"])
        )
    else:
        # Add fallback icon as a separate shape
        icon_left = Inches(IMAGE_AREA["left"] + IMAGE_AREA["width"]/2 - 0.5)
        icon_top = Inches(IMAGE_AREA["top"] + IMAGE_AREA["height"]/2 - 0.5)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(1), Inches(1))
        icon_tf = icon_box.text_frame
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("solution")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=72, alignment=PP_ALIGN.CENTER)

def create_custom_features_slide(prs: Presentation, content: Dict[str, Any], image_manager: ImageManager):
    """Create the customized key features slide."""
    slide_layout = prs.slide_layouts[5]  # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = content["title"]
    apply_text_formatting(
        title.text_frame, 
        font_name=FONTS["title"], 
        size=FONT_SIZES["title"], 
        bold=True
    )
    
    # Create table for features
    features = content["features"]
    rows, cols = len(features), 2
    
    left = Inches(CONTENT_AREA["left"])
    top = Inches(CONTENT_AREA["top"])
    width = Inches(CONTENT_AREA["width"] + IMAGE_AREA["width"])
    height = Inches(CONTENT_AREA["height"])
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths - adjusted for 16:9 format
    table.columns[0].width = Inches(0.5)  # Icon column
    table.columns[1].width = Inches(11.83)  # Feature text column - wider for 16:9
    
    # Add features to table
    for i, feature in enumerate(features):
        # Icon cell
        icon_cell = table.cell(i, 0)
        icon = match_icon_to_feature(feature)
        icon_cell.text = icon
        apply_text_formatting(icon_cell.text_frame, size=16, alignment=PP_ALIGN.CENTER)
        
        # Feature cell
        feature_cell = table.cell(i, 1)
        feature_cell.text = feature
        apply_text_formatting(feature_cell.text_frame)
        
        # Set row height
        table.rows[i].height = Inches(CONTENT_AREA["height"] / len(features))
    
    # Style table with no fill
    for cell in table.iter_cells():
        cell.fill.background()
    
    # Add custom image if available (cached or uploaded), positioned below table
    custom_image = image_manager.get_image_for_slide('features_slide')
    if custom_image:
        pic = slide.shapes.add_picture(
            custom_image,
            Inches(9.83),  # Further right for 16:9 format
            Inches(5.5),
            Inches(2.5),
            Inches(1.5)
        )

def create_custom_advantage_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any], image_manager: ImageManager):
    """Create the customized competitive advantage slide."""
    # Use blank layout to avoid unwanted placeholder elements
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title manually
    title_left = Inches(1.0)
    title_top = Inches(0.5)
    title_width = Inches(11.33)
    title_height = Inches(1.0)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = content["title"]
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting and allow natural wrapping
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(FONT_SIZES["title"])
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Ensure title wraps naturally if it's too long
    title_tf.word_wrap = True
    
    # Create custom text box for bullet points with better spacing
    left = Inches(CONTENT_AREA["left"])
    top = Inches(CONTENT_AREA["top"])
    width = Inches(CONTENT_AREA["width"])  # Use defined width to prevent overflow into image area
    height = Inches(CONTENT_AREA["height"])
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_right = Inches(0.1)  # Add right margin to prevent text overflow
    
    # Clear any existing text and add bullet points
    if tf.paragraphs:
        p = tf.paragraphs[0]
        p.text = ""
    else:
        p = tf.add_paragraph()
        
    # Add bullet points with proper formatting and bullet characters
    for i, bullet in enumerate(content["bullets"]):
        # Truncate bullet points if they're too long (prevents text overflow)
        bullet_text = bullet
        if len(bullet) > 120:  # Reasonable character limit for bullet points
            bullet_text = bullet[:117] + "..."
            
        if i == 0:
            p.text = f"• {bullet_text}"  # Add bullet character
        else:
            p = tf.add_paragraph()
            p.text = f"• {bullet_text}"  # Add bullet character
        
        # Apply bullet formatting
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        
        # Apply font formatting
        p.font.name = FONTS["body"]
        p.font.size = Pt(FONT_SIZES["body"])
        
        # Apply spacing for better readability
        p.space_before = Pt(6)  # Add space before each bullet
        p.space_after = Pt(6)   # Add space after each bullet
        
    apply_text_formatting(tf)
    
    # Add custom image directly to slide without using placeholder
    image_data = image_manager.get_image_for_slide('advantage_slide')
    
    if image_data:
        # Position the image on the right side of the 16:9 slide with proper dimensions
        pic = slide.shapes.add_picture(
            image_data,
            Inches(IMAGE_AREA["left"]),     # Positioned further right for 16:9 format
            Inches(IMAGE_AREA["top"]),
            Inches(IMAGE_AREA["width"]),    # Adjusted width for 16:9
            Inches(IMAGE_AREA["height"])
        )
    else:
        # Add fallback icon as a separate shape
        icon_left = Inches(IMAGE_AREA["left"] + IMAGE_AREA["width"]/2 - 0.5)
        icon_top = Inches(IMAGE_AREA["top"] + IMAGE_AREA["height"]/2 - 0.5)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(1), Inches(1))
        icon_tf = icon_box.text_frame
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("advantage")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=72, alignment=PP_ALIGN.CENTER)

def create_custom_audience_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any], image_manager: ImageManager):
    """Create the customized target audience slide."""
    # Use blank layout to avoid unwanted placeholder elements
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title manually
    title_left = Inches(1.0)
    title_top = Inches(0.5)
    title_width = Inches(11.33)
    title_height = Inches(1.0)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = content["title"]
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting and allow natural wrapping
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(FONT_SIZES["title"])
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Ensure title wraps naturally if it's too long
    title_tf.word_wrap = True
    
    # Create custom text box for paragraph content with better spacing
    left = Inches(CONTENT_AREA["left"])
    top = Inches(CONTENT_AREA["top"])
    width = Inches(CONTENT_AREA["width"])  # Use defined width to prevent overflow into image area
    height = Inches(CONTENT_AREA["height"])
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_right = Inches(0.1)  # Add right margin to prevent text overflow

    p = tf.paragraphs[0]
    # Truncate paragraph if it's too long to prevent text overflow
    paragraph_text = content["paragraph"]
    if len(paragraph_text) > 550:  # Reasonable limit for paragraphs
        paragraph_text = paragraph_text[:547] + "..."
    p.text = paragraph_text
    p.alignment = PP_ALIGN.LEFT
    
    # Format paragraph with proper spacing
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    p.line_spacing = 1.2  # Better line spacing
    
    apply_text_formatting(tf)
    
    # Add custom image directly to slide without using placeholder
    image_data = image_manager.get_image_for_slide('audience_slide')
    
    if image_data:
        # Position the image on the right side of the 16:9 slide with proper dimensions
        pic = slide.shapes.add_picture(
            image_data,
            Inches(IMAGE_AREA["left"]),     # Positioned further right for 16:9 format
            Inches(IMAGE_AREA["top"]),
            Inches(IMAGE_AREA["width"]),    # Adjusted width for 16:9
            Inches(IMAGE_AREA["height"])
        )
    else:
        # Add fallback icon as a separate shape
        icon_left = Inches(IMAGE_AREA["left"] + IMAGE_AREA["width"]/2 - 0.5)
        icon_top = Inches(IMAGE_AREA["top"] + IMAGE_AREA["height"]/2 - 0.5)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(1), Inches(1))
        icon_tf = icon_box.text_frame
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("audience")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=72, alignment=PP_ALIGN.CENTER)

def create_custom_cta_slide(prs: Presentation, content: Dict[str, Any], image_manager: ImageManager):
    """Create the customized call to action slide."""
    slide_layout = prs.slide_layouts[5]  # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Get Started"
    apply_text_formatting(
        title.text_frame, 
        font_name=FONTS["title"], 
        size=FONT_SIZES["title"], 
        bold=True
    )
    
    # CTA text (center of slide)
    left = Inches(1.0)  # Increased margin for better spacing
    top = Inches(2.5)
    width = Inches(11.33)  # Wider for 16:9 format
    height = Inches(1.5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = content["call_to_action"]
    p.alignment = PP_ALIGN.CENTER
    
    apply_text_formatting(
        tf, 
        font_name=FONTS["title"], 
        size=FONT_SIZES["call_to_action"], 
        bold=True,
        alignment=PP_ALIGN.CENTER
    )
    
    # Add custom image if available (cached or uploaded), positioned at bottom
    custom_image = image_manager.get_image_for_slide('cta_slide')
    if custom_image:
        pic = slide.shapes.add_picture(
            custom_image,
            Inches(5.16),  # Centered for 16:9 format
            Inches(4.5),
            Inches(3),
            Inches(2)
        )
    else:
        # Optional grey stripe at bottom
        left = Inches(0)
        top = Inches(6.5)
        width = Inches(13.33)  # Full width of 16:9 slide
        height = Inches(0.5)
        
        stripe_box = slide.shapes.add_textbox(left, top, width, height)
        stripe_box.fill.solid()
        stripe_box.fill.fore_color.rgb = RGBColor.from_string(COLORS["gray"])
        stripe_box.line.fill.background()
