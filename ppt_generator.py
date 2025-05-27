import io
import streamlit as st
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from image_fetcher import fetch_image_for_slide, get_slide_icon
from utils import FONTS, COLORS, FONT_SIZES, MARGINS, CONTENT_AREA, IMAGE_AREA, match_icon_to_feature
# Import additional slide creation functions
from ppt_generator_additions import create_market_slide, create_roadmap_slide, create_team_slide

def create_presentation(content: Dict[str, Any], filename: str, image_manager=None, custom_slide_order: List[str] = None) -> io.BytesIO:
    """
    Create a PowerPoint presentation from generated content with optional custom slide order.
    
    Args:
        content: Dictionary containing content for all slides
        filename: Name to give the file
        image_manager: Optional image manager for custom images
        custom_slide_order: Optional list specifying the order of slides
        
    Returns:
        BytesIO object containing the presentation
    """
    import streamlit as st
    import hashlib
    import json
    
    # Create a cache key based on the content and slide order
    cache_key = None
    try:
        # Create a deterministic representation of the inputs
        cache_input = {
            'content': content,
            'filename': filename,
            'custom_slide_order': custom_slide_order if custom_slide_order else []
        }
        # Hash the stringified content to create a cache key
        cache_key = hashlib.md5(json.dumps(cache_input, sort_keys=True).encode()).hexdigest()
        
        # Check if we have a cached presentation
        if 'ppt_cache' in st.session_state and cache_key in st.session_state.ppt_cache:
            print("Using cached presentation - no regeneration needed")
            # Get a fresh copy of the cached BytesIO
            cached_ppt = st.session_state.ppt_cache[cache_key]
            cached_ppt.seek(0)
            return cached_ppt
    except Exception as e:
        print(f"Cache key generation error: {e}")
        # Continue without caching if there's an error
    
    # Flag to indicate we're creating a new presentation (not from cache)
    # This will be used to determine whether to fetch new images
    is_new_presentation = True
    
    prs = Presentation()
    
    # Set slide size to widescreen (16:9) - 13.33 x 7.5 inches
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Initialize image cache for future customizations
    if 'original_images_cache' not in st.session_state:
        st.session_state.original_images_cache = {}
    
    # Track used images in this presentation to prevent duplicates
    used_images = set()
    
    # Store the flag in the presentation context
    presentation_context = {
        'metadata': content.get('metadata', {}),
        'is_new_presentation': is_new_presentation  # Add this flag
    }
    
    # Determine slide order
    if custom_slide_order:
        # Filter out any slide types that don't exist in the content
        slide_order = [slide_type for slide_type in custom_slide_order if slide_type in content]
        
        # For safety, add any missing slides that are in content but not in slide_order
        for slide_type in content:
            if slide_type.endswith('_slide') and slide_type not in slide_order:
                slide_order.append(slide_type)
    else:
        # Default order: title first, core slides, then optional slides, CTA last
        core_slides = ['title_slide', 'problem_slide', 'solution_slide', 'features_slide', 
                       'advantage_slide', 'audience_slide']
        optional_slides = ['market_slide', 'roadmap_slide', 'team_slide']
        cta_slide = ['cta_slide']
        
        slide_order = []
        
        # Add core slides that exist in content
        for slide_type in core_slides:
            if slide_type in content:
                slide_order.append(slide_type)
        
        # Add optional slides that exist in content
        for slide_type in optional_slides:
            if slide_type in content:
                slide_order.append(slide_type)
        
        # CTA slide always last if it exists
        if 'cta_slide' in content:
            slide_order.append('cta_slide')
    
    # Debug info for slide order
    print(f"Creating presentation with slides in order: {slide_order}")
    
    # Create slides in the specified order
    for slide_type in slide_order:
        if slide_type not in content:
            continue
            
        slide_data = content[slide_type]
        
        # Create the appropriate slide based on type
        if slide_type == 'title_slide':
            create_title_slide(prs, slide_data)
        elif slide_type == 'problem_slide':
            create_problem_slide(prs, slide_data, content, used_images)
        elif slide_type == 'solution_slide':
            create_solution_slide(prs, slide_data, content, used_images)
        elif slide_type == 'features_slide':
            create_features_slide(prs, slide_data)
        elif slide_type == 'advantage_slide':
            create_advantage_slide(prs, slide_data, content, used_images)
        elif slide_type == 'audience_slide':
            create_audience_slide(prs, slide_data, content, used_images)
        elif slide_type == 'market_slide':
            # Using the wrapper function instead
            create_market_slide_wrapper(prs, slide_data, content, used_images)
        elif slide_type == 'roadmap_slide':
            # Using the wrapper function instead
            create_roadmap_slide_wrapper(prs, slide_data, content, used_images)
        elif slide_type == 'team_slide':
            # Using the wrapper function instead
            create_team_slide_wrapper(prs, slide_data, content, used_images)
        elif slide_type == 'cta_slide':
            create_cta_slide(prs, slide_data)
      # Save to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    # Cache the presentation in session state
    try:
        if cache_key:
            # Initialize the cache if it doesn't exist
            if 'ppt_cache' not in st.session_state:
                st.session_state.ppt_cache = {}
            
            # Create a new BytesIO to store in cache (since the returned one will be consumed)
            cache_copy = io.BytesIO(output.getvalue())
            st.session_state.ppt_cache[cache_key] = cache_copy
            print("Presentation cached successfully")
    except Exception as e:
        print(f"Error caching presentation: {e}")
    
    return output

def apply_text_formatting(text_frame, font_name=FONTS["body"], size=FONT_SIZES["body"], 
                         bold=False, color=COLORS["black"], alignment=PP_ALIGN.LEFT):
    """Apply consistent text formatting to a text frame."""
    text_frame.word_wrap = True
    text_frame.auto_size = False  # Prevent auto-sizing to control overflow
    
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = alignment
        paragraph.line_spacing = 1.2  # Better line spacing for readability
        
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor.from_string(color)

def truncate_text_for_slide(text, max_chars=None, max_words=None):
    """Intelligently truncate text for slides to prevent overflow."""
    if not text:
        return ""
        
    # For very long text, apply stricter limits regardless of provided parameters
    if len(text) > 500:
        max_chars = min(max_chars or 400, 400)
        max_words = min(max_words or 75, 75)
        
    # If both limits are provided, use the most restrictive one
    if max_chars and max_words:
        # First check word limit
        words = text.split()
        if len(words) > max_words:
            truncated = " ".join(words[:max_words]) + "..."
            # Then check char limit on the word-truncated text
            if len(truncated) > max_chars:
                return truncated[:max_chars-3] + "..."
            return truncated
    
    # Check character limit
    if max_chars and len(text) > max_chars:
        return text[:max_chars-3] + "..."

    # Check word limit
    if max_words:
        words = text.split()
        if len(words) > max_words:
            return " ".join(words[:max_words]) + "..."
            
    return text

def create_title_slide(prs: Presentation, content: Dict[str, str]):
    """Create the title slide."""
    # Use blank layout for consistent control across all slides
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background to white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(COLORS["white"])
    
    # Title - centered and positioned at top third
    title_left = Inches(0.5)
    title_top = Inches(2.0)  # Move down from top
    title_width = Inches(12.33)
    title_height = Inches(2.0)  # Taller to handle multi-line titles
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    
    # Replace product name placeholder and truncate if needed
    product_name = content.get("product_name", "")
    clean_title = content["title"].replace("[Product Name]", product_name)
    # Limit title length to prevent overflow
    title_p.text = truncate_text_for_slide(clean_title, max_chars=60, max_words=10)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Apply bold, large title formatting with proper font
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(40)  # Larger size for title slide
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Subtitle - centered and positioned below title
    subtitle_left = Inches(0.5)
    subtitle_top = Inches(4.0)  # Position below title
    subtitle_width = Inches(12.33)
    subtitle_height = Inches(1.0)
    
    subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    subtitle_tf.auto_size = False
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_p.text = truncate_text_for_slide(content["subtitle"], max_chars=80)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Apply subtitle formatting
    for run in subtitle_p.runs:
        run.font.name = FONTS["body"]
        run.font.size = Pt(28)  # Larger for better visibility
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])

def create_problem_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the problem statement slide with improved text handling."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Add title with clear positioning
    title_left = Inches(0.5)  # Consistent margin
    title_top = Inches(0.5)   # Consistent position from top
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    
    # Replace any placeholder text and truncate if needed
    clean_title = content["title"].replace("[Product Name]", product_name)
    title_p.text = truncate_text_for_slide(clean_title, max_chars=60)
    title_p.alignment = PP_ALIGN.LEFT  # Left aligned for content slides
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)  # Consistent title size
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Reduce width to ensure content doesn't overflow when image is present
    content_width = Inches(6.8)  # Narrower to avoid potential overflow
    
    # Create text box for bullet points with better margins
    left = Inches(0.5)
    top = Inches(1.7)
    width = content_width
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    # Clear any existing text
    if tf.paragraphs:
        p = tf.paragraphs[0]
        p.text = ""
    else:
        p = tf.add_paragraph()
        
    # Handle different field names from slide content generators
    bullets = []
    if "pain_points" in content:
        bullets = content["pain_points"]
    elif "bullets" in content:
        bullets = content["bullets"]
    elif "differentiators" in content:
        bullets = [item["point"] if isinstance(item, dict) else str(item) for item in content["differentiators"]]
    
    # Keep max 3-4 bullets to prevent crowding and ensure spacing
    max_bullets = 4
    bullets = bullets[:max_bullets]
    
    # More aggressive truncation for bullet text
    max_bullet_chars = 75  # Reduced character limit
    max_bullet_words = 15  # Reduced word limit
    
    # Apply better spacing between bullets based on bullet count
    bullet_count = len(bullets)
    if bullet_count <= 2:
        bullet_spacing = Pt(20)  # More space for few bullets
    elif bullet_count == 3:
        bullet_spacing = Pt(15)
    else:
        bullet_spacing = Pt(12)  # Less space but still adequate for 4+ bullets
    
    for i, bullet in enumerate(bullets):
        # Replace product name placeholders
        clean_bullet = bullet.replace("[Product Name]", product_name) if product_name else bullet
        
        # Truncate bullet points with stricter limits
        bullet_text = truncate_text_for_slide(clean_bullet, max_chars=max_bullet_chars, max_words=max_bullet_words)
            
        if i == 0:
            p.text = f"• {bullet_text}"
        else:
            p = tf.add_paragraph()
            p.text = f"• {bullet_text}"
        
        # Apply improved bullet formatting
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_before = bullet_spacing
        p.space_after = bullet_spacing
        # Adjust line spacing based on bullet count
        p.line_spacing = 1.1 if bullet_count > 3 else 1.2
        
        # Apply font formatting directly to runs
        for run in p.runs:
            run.font.name = FONTS["body"]
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Add image with proper positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("problem", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        # Position image on right side - slightly adjusted for better balance
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),   # Move slightly left to avoid edge
            Inches(2.0),   # Centered vertically
            Inches(4.5),   # Consistent width
            Inches(4.0)    # Consistent height
        )
    else:
        # Add fallback icon with better positioning
        icon_left = Inches(9.5)
        icon_top = Inches(3.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_info = get_slide_icon("problem")
        icon_p.text = icon_info["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=64, alignment=PP_ALIGN.CENTER)

def create_solution_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the solution overview slide with improved text containment."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    
    # Replace any placeholder text and truncate
    clean_title = content["title"].replace("[Product Name]", product_name)
    title_p.text = truncate_text_for_slide(clean_title, max_chars=60)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Create custom text box for paragraph content - narrower width for better text fitting
    content_width = Inches(6.8)  # Reduced width to prevent overflow
    
    left = Inches(0.5)
    top = Inches(1.7)
    width = content_width
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    p = tf.paragraphs[0]
    
    # Handle different field names for content
    paragraph_text = ""
    if "paragraph" in content:
        paragraph_text = content["paragraph"]
    elif "description" in content:
        paragraph_text = content["description"]
    elif "value_proposition" in content:
        paragraph_text = content["value_proposition"]
    
    # Replace product name placeholders
    clean_paragraph = paragraph_text.replace("[Product Name]", product_name) if product_name else paragraph_text
    
    # More aggressive truncation for better text fitting
    p.text = truncate_text_for_slide(clean_paragraph, max_chars=350, max_words=70)
    
    # Format paragraph with proper spacing
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    p.line_spacing = 1.2
    
    # Apply paragraph formatting directly to runs
    for run in p.runs:
        run.font.name = FONTS["body"]
        run.font.size = Pt(18)  # Slightly smaller for better fit
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Add image with consistent positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("solution", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        # Position image on right side with consistent sizing
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),  # Adjusted position
            Inches(2.0),
            Inches(4.5),
            Inches(4.0)
        )
    else:
        # Add fallback icon
        icon_left = Inches(9.5)
        icon_top = Inches(3.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("solution")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=64, alignment=PP_ALIGN.CENTER)

def create_features_slide(prs: Presentation, content: Dict[str, Any]):
    """Create the key features slide with improved spacing."""
    slide_layout = prs.slide_layouts[6]  # Blank layout for consistency
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = truncate_text_for_slide(content["title"].replace("[Product Name]", ""), max_chars=60)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Get features list
    features = content["features"]
    max_features = 5  # Limit to 5 features for better spacing
    features = features[:max_features]
    
    # Calculate spacing based on number of features to prevent overlap
    if len(features) <= 3:
        feature_height = Inches(1.5)  # More space for few features
    else:
        feature_height = Inches(1.2)  # Less space but still adequate for more features
    
    # Create a more visually appealing list with proper spacing
    left_margin = Inches(1.0)
    top_start = Inches(1.7)
    icon_width = Inches(0.8)
    text_width = Inches(11.0)
    
    for i, feature in enumerate(features):
        # Handle different feature formats
        if isinstance(feature, dict):
            feature_text = feature.get('feature', feature.get('name', feature.get('title', str(feature))))
            icon_text = feature_text
        else:
            feature_text = str(feature)
            icon_text = feature_text
        
        # Truncate feature text appropriately
        feature_text = truncate_text_for_slide(feature_text, max_chars=100, max_words=20)

        # Calculate vertical position with dynamic spacing
        top_position = top_start + (i * feature_height)
        
        # Add icon
        icon_box = slide.shapes.add_textbox(left_margin, top_position, icon_width, feature_height)
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = match_icon_to_feature(icon_text)
        icon_p.alignment = PP_ALIGN.CENTER
        
        # Apply icon formatting
        for run in icon_p.runs:
            run.font.name = FONTS["body"]
            run.font.size = Pt(24)
            run.font.bold = True
        
        # Add feature text
        text_left = left_margin + icon_width + Inches(0.3)
        text_box = slide.shapes.add_textbox(text_left, top_position, text_width, feature_height)
        text_tf = text_box.text_frame
        text_tf.word_wrap = True
        text_tf.auto_size = False
        text_p = text_tf.paragraphs[0]
        text_p.text = feature_text
        text_p.alignment = PP_ALIGN.LEFT
        
        # Set line spacing to prevent overlap
        text_p.space_before = Pt(3)
        text_p.space_after = Pt(3)
        text_p.line_spacing = 1.1
        
        # Apply text formatting
        for run in text_p.runs:
            run.font.name = FONTS["body"]
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor.from_string(COLORS["black"])

def create_advantage_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the competitive advantage slide with improved text containment."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = content["title"].replace("[Product Name]", product_name)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Create text box for bullets - narrower for better text fitting
    content_width = Inches(6.8)  # Reduced width to prevent overflow
    
    left = Inches(0.5)
    top = Inches(1.7)
    width = content_width
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    # Clear any existing text
    if tf.paragraphs:
        p = tf.paragraphs[0]
        p.text = ""
    else:
        p = tf.add_paragraph()
        
    # Handle different field names for bulleted content
    bullets = []
    if "differentiators" in content:
        # Handle various possible dictionary structures safely
        differentiators = content["differentiators"]
        for item in differentiators:
            if isinstance(item, dict):
                # Try various possible key names for the point
                if "point" in item:
                    bullets.append(item["point"])
                elif "text" in item:
                    bullets.append(item["text"])
                elif "description" in item:
                    bullets.append(item["description"])
                elif "name" in item:
                    bullets.append(item["name"])
                elif len(item) > 0:
                    # Just take the first value if we can't find a known key
                    bullets.append(next(iter(item.values())))
                else:
                    # Empty dict, add placeholder
                    bullets.append("Competitive advantage point")
            else:
                # Handle non-dict items (strings, etc.)
                bullets.append(str(item))
    elif "bullets" in content:
        bullets = content["bullets"]
    elif "pain_points" in content:
        bullets = content["pain_points"]
    
    # If we still have no bullets, add a placeholder
    if not bullets:
        bullets = ["This product offers significant competitive advantages"]
    
    # Limit bullets for better readability
    max_bullets = 4
    bullets = bullets[:max_bullets]
    
    # More aggressive truncation for better text fitting
    max_bullet_chars = 70  # Reduced character limit
    
    for i, bullet in enumerate(bullets):
        # Replace any product name placeholders
        clean_bullet = bullet.replace("[Product Name]", product_name) if product_name and isinstance(bullet, str) else str(bullet)
        
        # Truncate with stricter limits
        bullet_text = truncate_text_for_slide(clean_bullet, max_chars=max_bullet_chars, max_words=14)
        
        if i == 0:
            p.text = f"• {bullet_text}"
        else:
            p = tf.add_paragraph()
            p.text = f"• {bullet_text}"
        
        # Apply bullet formatting
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.font.name = FONTS["body"]
        p.font.size = Pt(18)  # Slightly smaller for better fit
        p.space_before = Pt(10)
        p.space_after = Pt(10)
        p.line_spacing = 1.1  # Tighter line spacing
    
    # Add image with consistent positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("advantage", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        # Position image with adjusted dimensions
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),  # Adjusted position
            Inches(2.0),
            Inches(4.5),
            Inches(4.0)
        )
    else:
        # Add fallback icon
        icon_left = Inches(9.5)
        icon_top = Inches(3.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("advantage")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=64, alignment=PP_ALIGN.CENTER)

def create_audience_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the target audience slide with improved text containment."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = content["title"].replace("[Product Name]", product_name)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Create text box for content with consistent positioning - narrower for better text fitting
    content_width = Inches(6.8)  # Reduced width to prevent overflow
    
    left = Inches(0.5)
    top = Inches(1.7)
    width = content_width
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    p = tf.paragraphs[0]
    
    # Handle different field names for text content
    paragraph_text = ""
    if "paragraph" in content:
        paragraph_text = content["paragraph"]
    elif "description" in content:
        paragraph_text = content["description"]
    elif "content" in content:
        paragraph_text = content["content"]
    
    # Replace product name placeholders
    clean_paragraph = paragraph_text.replace("[Product Name]", product_name) if product_name else paragraph_text
    
    # More aggressive truncation for better text fitting
    p.text = truncate_text_for_slide(clean_paragraph, max_chars=350, max_words=70)
    
    # Format paragraph with proper spacing
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    p.line_spacing = 1.2
    p.font.name = FONTS["body"]
    p.font.size = Pt(18)  # Slightly smaller for better fit
    
    # Add image with consistent positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("audience", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        # Position image with consistent dimensions
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),  # Adjusted position
            Inches(2.0),
            Inches(4.5),
            Inches(4.0)
        )
    else:
        # Add fallback icon
        icon_left = Inches(9.5)
        icon_top = Inches(3.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("audience")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=64, alignment=PP_ALIGN.CENTER)

def create_cta_slide(prs: Presentation, content: Dict[str, Any]):
    """Create the call to action slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Main centered title
    title_left = Inches(0.5)
    title_top = Inches(1.5)  # Higher placement
    title_width = Inches(12.33)
    title_height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = content.get("title", "Get Started Today")
    title_p.alignment = PP_ALIGN.CENTER
    
    # Apply title formatting - larger and bolder
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(36)  # Larger for emphasis
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # CTA text - centered and prominent
    cta_left = Inches(1.0)
    cta_top = Inches(3.0)  # Positioned below title
    cta_width = Inches(11.33)
    cta_height = Inches(2.0)  # Taller for multi-line CTAs
    
    cta_box = slide.shapes.add_textbox(cta_left, cta_top, cta_width, cta_height)
    cta_tf = cta_box.text_frame
    cta_tf.word_wrap = True
    cta_tf.auto_size = False
    
    cta_p = cta_tf.paragraphs[0]
    cta_text = content.get("cta_text", content.get("call_to_action", "Contact us today to learn more!"))
    
    # Don't truncate CTAs - they're often short and important
    cta_p.text = cta_text
    cta_p.alignment = PP_ALIGN.CENTER
    
    # Apply CTA formatting - prominent and attention-grabbing
    for run in cta_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Add any supporting bullets/text if available
    if "bullets" in content and content["bullets"]:
        bullet_top = Inches(5.0)  # Position below CTA
        bullet_height = Inches(1.5)
        
        bullet_box = slide.shapes.add_textbox(cta_left, bullet_top, cta_width, bullet_height)
        bullet_tf = bullet_box.text_frame
        bullet_tf.word_wrap = True
        
        # Add up to 2 bullets only
        max_bullets = min(2, len(content["bullets"]))
        for i in range(max_bullets):
            if i == 0:
                bullet_p = bullet_tf.paragraphs[0]
            else:
                bullet_p = bullet_tf.add_paragraph()
                
            bullet_p.text = f"• {content['bullets'][i]}"
            bullet_p.alignment = PP_ALIGN.CENTER
            
            # Format bullet text
            for run in bullet_p.runs:
                run.font.name = FONTS["body"]
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Optional footer bar
    footer_left = Inches(0)
    footer_top = Inches(6.5)
    footer_width = Inches(13.33)
    footer_height = Inches(1.0)
    
    footer_box = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
    footer_box.fill.solid()
    footer_box.fill.fore_color.rgb = RGBColor.from_string(COLORS["gray"])
    footer_box.line.fill.background()

def fetch_image_with_cache(slide_type: str, context: Dict[str, Any], used_images: set, use_placeholders: bool = False):
    """Fetch an image with caching and deduplication to prevent repeating the same image."""
    import streamlit as st
    
    # Check if we're creating a new presentation or using a cached one
    is_new_presentation = context.get('is_new_presentation', True)
    
    # Check if we already have this image in cache
    if 'original_images_cache' in st.session_state and f'{slide_type}_slide' in st.session_state.original_images_cache:
        # Get the cached image data
        cached_data = st.session_state.original_images_cache[f'{slide_type}_slide']
        
        # Create a hash of the image data to track unique images
        import hashlib
        img_hash = hashlib.md5(cached_data).hexdigest()
        
        # Check if this image has already been used in this presentation
        if img_hash not in used_images:
            # Mark image as used
            used_images.add(img_hash)
            # Return the cached image data
            return io.BytesIO(cached_data)
        elif not is_new_presentation:
            # If we're not creating a new presentation (using cache),
            # just return the cached image even if it's used before
            # This prevents unnecessary API calls
            return io.BytesIO(cached_data)
        else:
            # Image has been used before in a new presentation, fetch a different one
            return fetch_unique_image(slide_type, context, used_images, use_placeholders)
    
    # No cached image, fetch a new one only if creating a new presentation
    if is_new_presentation:
        return fetch_unique_image(slide_type, context, used_images, use_placeholders)
    else:
        # For cached presentations, don't fetch new images if not in cache
        return None

def fetch_unique_image(slide_type: str, context: Dict[str, Any], used_images: set, use_placeholders: bool = False):
    """Fetch a unique image that hasn't been used in the presentation yet."""
    # Don't fetch new images if we're not creating a new presentation
    if not context.get('is_new_presentation', True):
        return None
        
    # Try up to 3 times to get a unique image
    for _ in range(3):
        # Fetch a new image
        image_data = fetch_image_for_slide(slide_type, context, use_placeholders)
        if not image_data:
            return None
            
        # Check if unique
        image_data.seek(0)
        img_bytes = image_data.read()
        image_data.seek(0)
        
        import hashlib
        img_hash = hashlib.md5(img_bytes).hexdigest()
        
        if img_hash not in used_images:
            # Mark as used and cache
            used_images.add(img_hash)
            
            # Cache the image
            import streamlit as st
            if 'original_images_cache' in st.session_state:
                st.session_state.original_images_cache[f'{slide_type}_slide'] = img_bytes
            
            return image_data
    
    # If we couldn't get a unique image after 3 tries, use a generic one
    return None

def get_search_terms_for_slide(slide_type: str, context: Dict[str, Any]) -> List[str]:
    """Generate specific search terms for each slide type based on content."""
    search_terms = []
    product_name = context.get("metadata", {}).get("product_name", "") if context else ""
    
    if slide_type == "problem":
        # Get keywords from problem statement
        if "problem_slide" in context and "pain_points" in context["problem_slide"]:
            pain_points = context["problem_slide"]["pain_points"]
            for point in pain_points[:2]:  # Use first two pain points for keywords
                # Extract key nouns from the pain point
                words = [w for w in point.split() if len(w) > 4 and w.lower() not in 
                        ["their", "there", "these", "those", "about", "because", "through"]]
                search_terms.extend(words[:2])  # Add up to 2 keywords per pain point
        search_terms.extend(["problem", "challenge", "frustration"])
        
    elif slide_type == "solution":
        # Get keywords from solution description
        if "solution_slide" in context and "paragraph" in context["solution_slide"]:
            solution_text = context["solution_slide"]["paragraph"]
            # Extract key nouns from solution
            words = [w for w in solution_text.split() if len(w) > 5 and w.lower() not in 
                    ["their", "there", "these", "those", "about", "because", "through"]]
            search_terms.extend(words[:3])  # Add up to 3 keywords
        search_terms.extend([product_name, "solution", "innovation"])
        
    elif slide_type == "advantage":
        # Get keywords from advantages
        if "advantage_slide" in context and "differentiators" in context["advantage_slide"]:
            advantages = context["advantage_slide"]["differentiators"]
            for adv in advantages[:2]:  # Use first two advantages for keywords
                if isinstance(adv, str):
                    words = adv.split()
                    search_terms.extend([w for w in words if len(w) > 5][:1])
        search_terms.extend(["advantage", "unique", "competitive edge"])
        
    elif slide_type == "audience":
        # Get keywords from audience description
        if "audience_slide" in context and "paragraph" in context["audience_slide"]:
            audience_text = context["audience_slide"]["paragraph"]
            # Extract potential demographic keywords
            demographics = ["business", "enterprise", "consumer", "professional", 
                           "corporate", "small business", "tech", "healthcare", 
                           "finance", "retail", "education"]
            for demo in demographics:
                if demo in audience_text.lower():
                    search_terms.append(demo)
        search_terms.extend(["target audience", "customer", "user", "client"])
        
    elif slide_type == "market":
        # Use market-related terms
        if "market_slide" in context:
            market_desc = context["market_slide"].get("description", "")
            # Look for industry keywords
            industries = ["tech", "healthcare", "finance", "retail", "education", 
                         "manufacturing", "automotive", "energy", "telecom"]
            for ind in industries:
                if ind in market_desc.lower():
                    search_terms.append(ind)
        search_terms.extend(["market growth", "market opportunity", "business chart", "market analysis"])
    
    # Ensure we have at least some search terms
    if not search_terms:
        if product_name:
            search_terms = [product_name, slide_type]
        else:
            search_terms = [slide_type, "business"]
    
    # Clean up search terms - remove special characters and short words
    cleaned_terms = []
    for term in search_terms:
        if isinstance(term, str):
            # Remove special characters
            cleaned = ''.join(c for c in term if c.isalnum() or c.isspace())
            # Only include if substantive
            if cleaned and len(cleaned) > 3:
                cleaned_terms.append(cleaned)
    
    # Return unique search terms, with an upper limit
    unique_terms = list(set(cleaned_terms))
    return unique_terms[:5]  # Limit to 5 terms for better focus

# Improved text handling functions

def fit_text_to_box(text_frame, text, font_name=FONTS["body"], size=FONT_SIZES["body"], 
                   bold=False, color=COLORS["black"], alignment=PP_ALIGN.LEFT):
    """Format and fit text properly within a text box, adjusting size if needed."""
    text_frame.word_wrap = True
    text_frame.auto_size = False
    
    p = text_frame.paragraphs[0]
    p.alignment = alignment
    p.text = text
    
    for run in p.runs:
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)
    
    # Check if text might be too long (rough estimation)
    if len(text) > 300 and size > 16:
        # Reduce font size for very long text
        for run in p.runs:
            run.font.size = Pt(size - 2)

def create_problem_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the problem statement slide with improved text handling."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Add title with clear positioning
    title_left = Inches(0.5)  # Consistent margin
    title_top = Inches(0.5)   # Consistent position from top
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    
    # Replace any placeholder text and truncate if needed
    clean_title = content["title"].replace("[Product Name]", product_name)
    title_p.text = truncate_text_for_slide(clean_title, max_chars=60)
    title_p.alignment = PP_ALIGN.LEFT  # Left aligned for content slides
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)  # Consistent title size
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Reduce width to ensure content doesn't overflow when image is present
    content_width = Inches(6.8)  # Narrower to avoid potential overflow
    
    # Create text box for bullet points with better margins
    left = Inches(0.5)
    top = Inches(1.7)
    width = content_width
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    # Clear any existing text
    if tf.paragraphs:
        p = tf.paragraphs[0]
        p.text = ""
    else:
        p = tf.add_paragraph()
        
    # Handle different field names from slide content generators
    bullets = []
    if "pain_points" in content:
        bullets = content["pain_points"]
    elif "bullets" in content:
        bullets = content["bullets"]
    elif "differentiators" in content:
        bullets = [item["point"] if isinstance(item, dict) else str(item) for item in content["differentiators"]]
    
    # Keep max 3-4 bullets to prevent crowding and ensure spacing
    max_bullets = 4
    bullets = bullets[:max_bullets]
    
    # More aggressive truncation for bullet text
    max_bullet_chars = 75  # Reduced character limit
    max_bullet_words = 15  # Reduced word limit
    
    # Apply better spacing between bullets based on bullet count
    bullet_count = len(bullets)
    if bullet_count <= 2:
        bullet_spacing = Pt(20)  # More space for few bullets
    elif bullet_count == 3:
        bullet_spacing = Pt(15)
    else:
        bullet_spacing = Pt(12)  # Less space but still adequate for 4+ bullets
    
    for i, bullet in enumerate(bullets):
        # Replace product name placeholders
        clean_bullet = bullet.replace("[Product Name]", product_name) if product_name else bullet
        
        # Truncate bullet points with stricter limits
        bullet_text = truncate_text_for_slide(clean_bullet, max_chars=max_bullet_chars, max_words=max_bullet_words)
            
        if i == 0:
            p.text = f"• {bullet_text}"
        else:
            p = tf.add_paragraph()
            p.text = f"• {bullet_text}"
        
        # Apply improved bullet formatting
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_before = bullet_spacing
        p.space_after = bullet_spacing
        # Adjust line spacing based on bullet count
        p.line_spacing = 1.1 if bullet_count > 3 else 1.2
        
        # Apply font formatting directly to runs
        for run in p.runs:
            run.font.name = FONTS["body"]
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Add image with proper positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("problem", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        # Position image on right side - slightly adjusted for better balance
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),   # Move slightly left to avoid edge
            Inches(2.0),   # Centered vertically
            Inches(4.5),   # Consistent width
            Inches(4.0)    # Consistent height
        )
    else:
        # Add fallback icon with better positioning
        icon_left = Inches(9.5)
        icon_top = Inches(3.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_info = get_slide_icon("problem")
        icon_p.text = icon_info["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=64, alignment=PP_ALIGN.CENTER)

def create_solution_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the solution overview slide with improved text containment."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    
    # Replace any placeholder text and truncate
    clean_title = content["title"].replace("[Product Name]", product_name)
    title_p.text = truncate_text_for_slide(clean_title, max_chars=60)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Create custom text box for paragraph content - narrower width for better text fitting
    content_width = Inches(6.8)  # Reduced width to prevent overflow
    
    left = Inches(0.5)
    top = Inches(1.7)
    width = content_width
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    p = tf.paragraphs[0]
    
    # Handle different field names for content
    paragraph_text = ""
    if "paragraph" in content:
        paragraph_text = content["paragraph"]
    elif "description" in content:
        paragraph_text = content["description"]
    elif "value_proposition" in content:
        paragraph_text = content["value_proposition"]
    
    # Replace product name placeholders
    clean_paragraph = paragraph_text.replace("[Product Name]", product_name) if product_name else paragraph_text
    
    # More aggressive truncation for better text fitting
    p.text = truncate_text_for_slide(clean_paragraph, max_chars=350, max_words=70)
    
    # Format paragraph with proper spacing
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    p.line_spacing = 1.2
    
    # Apply paragraph formatting directly to runs
    for run in p.runs:
        run.font.name = FONTS["body"]
        run.font.size = Pt(18)  # Slightly smaller for better fit
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Add image with consistent positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("solution", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        # Position image on right side with consistent sizing
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),  # Adjusted position
            Inches(2.0),
            Inches(4.5),
            Inches(4.0)
        )
    else:
        # Add fallback icon
        icon_left = Inches(9.5)
        icon_top = Inches(3.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("solution")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=64, alignment=PP_ALIGN.CENTER)

def create_features_slide(prs: Presentation, content: Dict[str, Any]):
    """Create the key features slide with improved spacing."""
    slide_layout = prs.slide_layouts[6]  # Blank layout for consistency
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = truncate_text_for_slide(content["title"].replace("[Product Name]", ""), max_chars=60)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Get features list
    features = content["features"]
    max_features = 5  # Limit to 5 features for better spacing
    features = features[:max_features]
    
    # Calculate spacing based on number of features to prevent overlap
    if len(features) <= 3:
        feature_height = Inches(1.5)  # More space for few features
    else:
        feature_height = Inches(1.2)  # Less space but still adequate for more features
    
    # Create a more visually appealing list with proper spacing
    left_margin = Inches(1.0)
    top_start = Inches(1.7)
    icon_width = Inches(0.8)
    text_width = Inches(11.0)
    
    for i, feature in enumerate(features):
        # Handle different feature formats
        if isinstance(feature, dict):
            feature_text = feature.get('feature', feature.get('name', feature.get('title', str(feature))))
            icon_text = feature_text
        else:
            feature_text = str(feature)
            icon_text = feature_text
        
        # Truncate feature text appropriately
        feature_text = truncate_text_for_slide(feature_text, max_chars=100, max_words=20)

        # Calculate vertical position with dynamic spacing
        top_position = top_start + (i * feature_height)
        
        # Add icon
        icon_box = slide.shapes.add_textbox(left_margin, top_position, icon_width, feature_height)
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = match_icon_to_feature(icon_text)
        icon_p.alignment = PP_ALIGN.CENTER
        
        # Apply icon formatting
        for run in icon_p.runs:
            run.font.name = FONTS["body"]
            run.font.size = Pt(24)
            run.font.bold = True
        
        # Add feature text
        text_left = left_margin + icon_width + Inches(0.3)
        text_box = slide.shapes.add_textbox(text_left, top_position, text_width, feature_height)
        text_tf = text_box.text_frame
        text_tf.word_wrap = True
        text_tf.auto_size = False
        text_p = text_tf.paragraphs[0]
        text_p.text = feature_text
        text_p.alignment = PP_ALIGN.LEFT
        
        # Set line spacing to prevent overlap
        text_p.space_before = Pt(3)
        text_p.space_after = Pt(3)
        text_p.line_spacing = 1.1
        
        # Apply text formatting
        for run in text_p.runs:
            run.font.name = FONTS["body"]
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor.from_string(COLORS["black"])

def create_advantage_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the competitive advantage slide with improved text containment."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = content["title"].replace("[Product Name]", product_name)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Create text box for bullets - narrower for better text fitting
    content_width = Inches(6.8)  # Reduced width to prevent overflow
    
    left = Inches(0.5)
    top = Inches(1.7)
    width = content_width
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    # Clear any existing text
    if tf.paragraphs:
        p = tf.paragraphs[0]
        p.text = ""
    else:
        p = tf.add_paragraph()
        
    # Handle different field names for bulleted content
    bullets = []
    if "differentiators" in content:
        # Handle various possible dictionary structures safely
        differentiators = content["differentiators"]
        for item in differentiators:
            if isinstance(item, dict):
                # Try various possible key names for the point
                if "point" in item:
                    bullets.append(item["point"])
                elif "text" in item:
                    bullets.append(item["text"])
                elif "description" in item:
                    bullets.append(item["description"])
                elif "name" in item:
                    bullets.append(item["name"])
                elif len(item) > 0:
                    # Just take the first value if we can't find a known key
                    bullets.append(next(iter(item.values())))
                else:
                    # Empty dict, add placeholder
                    bullets.append("Competitive advantage point")
            else:
                # Handle non-dict items (strings, etc.)
                bullets.append(str(item))
    elif "bullets" in content:
        bullets = content["bullets"]
    elif "pain_points" in content:
        bullets = content["pain_points"]
    
    # If we still have no bullets, add a placeholder
    if not bullets:
        bullets = ["This product offers significant competitive advantages"]
    
    # Limit bullets for better readability
    max_bullets = 4
    bullets = bullets[:max_bullets]
    
    # More aggressive truncation for better text fitting
    max_bullet_chars = 70  # Reduced character limit
    
    for i, bullet in enumerate(bullets):
        # Replace any product name placeholders
        clean_bullet = bullet.replace("[Product Name]", product_name) if product_name and isinstance(bullet, str) else str(bullet)
        
        # Truncate with stricter limits
        bullet_text = truncate_text_for_slide(clean_bullet, max_chars=max_bullet_chars, max_words=14)
        
        if i == 0:
            p.text = f"• {bullet_text}"
        else:
            p = tf.add_paragraph()
            p.text = f"• {bullet_text}"
        
        # Apply bullet formatting
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.font.name = FONTS["body"]
        p.font.size = Pt(18)  # Slightly smaller for better fit
        p.space_before = Pt(10)
        p.space_after = Pt(10)
        p.line_spacing = 1.1  # Tighter line spacing
    
    # Add image with consistent positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("advantage", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        # Position image with adjusted dimensions
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),  # Adjusted position
            Inches(2.0),
            Inches(4.5),
            Inches(4.0)
        )
    else:
        # Add fallback icon
        icon_left = Inches(9.5)
        icon_top = Inches(3.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("advantage")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=64, alignment=PP_ALIGN.CENTER)

def create_audience_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the target audience slide with improved text containment."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = content["title"].replace("[Product Name]", product_name)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Create text box for content with consistent positioning - narrower for better text fitting
    content_width = Inches(6.8)  # Reduced width to prevent overflow
    
    left = Inches(0.5)
    top = Inches(1.7)
    width = content_width
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    p = tf.paragraphs[0]
    
    # Handle different field names for text content
    paragraph_text = ""
    if "paragraph" in content:
        paragraph_text = content["paragraph"]
    elif "description" in content:
        paragraph_text = content["description"]
    elif "content" in content:
        paragraph_text = content["content"]
    
    # Replace product name placeholders
    clean_paragraph = paragraph_text.replace("[Product Name]", product_name) if product_name else paragraph_text
    
    # More aggressive truncation for better text fitting
    p.text = truncate_text_for_slide(clean_paragraph, max_chars=350, max_words=70)
    
    # Format paragraph with proper spacing
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    p.line_spacing = 1.2
    p.font.name = FONTS["body"]
    p.font.size = Pt(18)  # Slightly smaller for better fit
    
    # Add image with consistent positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("audience", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        # Position image with consistent dimensions
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),  # Adjusted position
            Inches(2.0),
            Inches(4.5),
            Inches(4.0)
        )
    else:
        # Add fallback icon
        icon_left = Inches(9.5)
        icon_top = Inches(3.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("audience")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=64, alignment=PP_ALIGN.CENTER)

def create_cta_slide(prs: Presentation, content: Dict[str, Any]):
    """Create the call to action slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Main centered title
    title_left = Inches(0.5)
    title_top = Inches(1.5)  # Higher placement
    title_width = Inches(12.33)
    title_height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = content.get("title", "Get Started Today")
    title_p.alignment = PP_ALIGN.CENTER
    
    # Apply title formatting - larger and bolder
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(36)  # Larger for emphasis
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # CTA text - centered and prominent
    cta_left = Inches(1.0)
    cta_top = Inches(3.0)  # Positioned below title
    cta_width = Inches(11.33)
    cta_height = Inches(2.0)  # Taller for multi-line CTAs
    
    cta_box = slide.shapes.add_textbox(cta_left, cta_top, cta_width, cta_height)
    cta_tf = cta_box.text_frame
    cta_tf.word_wrap = True
    cta_tf.auto_size = False
    
    cta_p = cta_tf.paragraphs[0]
    cta_text = content.get("cta_text", content.get("call_to_action", "Contact us today to learn more!"))
    
    # Don't truncate CTAs - they're often short and important
    cta_p.text = cta_text
    cta_p.alignment = PP_ALIGN.CENTER
    
    # Apply CTA formatting - prominent and attention-grabbing
    for run in cta_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Add any supporting bullets/text if available
    if "bullets" in content and content["bullets"]:
        bullet_top = Inches(5.0)  # Position below CTA
        bullet_height = Inches(1.5)
        
        bullet_box = slide.shapes.add_textbox(cta_left, bullet_top, cta_width, bullet_height)
        bullet_tf = bullet_box.text_frame
        bullet_tf.word_wrap = True
        
        # Add up to 2 bullets only
        max_bullets = min(2, len(content["bullets"]))
        for i in range(max_bullets):
            if i == 0:
                bullet_p = bullet_tf.paragraphs[0]
            else:
                bullet_p = bullet_tf.add_paragraph()
                
            bullet_p.text = f"• {content['bullets'][i]}"
            bullet_p.alignment = PP_ALIGN.CENTER
            
            # Format bullet text
            for run in bullet_p.runs:
                run.font.name = FONTS["body"]
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Optional footer bar
    footer_left = Inches(0)
    footer_top = Inches(6.5)
    footer_width = Inches(13.33)
    footer_height = Inches(1.0)
    
    footer_box = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
    footer_box.fill.solid()
    footer_box.fill.fore_color.rgb = RGBColor.from_string(COLORS["gray"])
    footer_box.line.fill.background()

def fetch_image_with_cache(slide_type: str, context: Dict[str, Any], used_images: set, use_placeholders: bool = False):
    """Fetch an image with caching and deduplication to prevent repeating the same image."""
    import streamlit as st
    
    # Check if we're creating a new presentation or using a cached one
    is_new_presentation = context.get('is_new_presentation', True)
    
    # Check if we already have this image in cache
    if 'original_images_cache' in st.session_state and f'{slide_type}_slide' in st.session_state.original_images_cache:
        # Get the cached image data
        cached_data = st.session_state.original_images_cache[f'{slide_type}_slide']
        
        # Create a hash of the image data to track unique images
        import hashlib
        img_hash = hashlib.md5(cached_data).hexdigest()
        
        # Check if this image has already been used in this presentation
        if img_hash not in used_images:
            # Mark image as used
            used_images.add(img_hash)
            # Return the cached image data
            return io.BytesIO(cached_data)
        elif not is_new_presentation:
            # If we're not creating a new presentation (using cache),
            # just return the cached image even if it's used before
            # This prevents unnecessary API calls
            return io.BytesIO(cached_data)
        else:
            # Image has been used before in a new presentation, fetch a different one
            return fetch_unique_image(slide_type, context, used_images, use_placeholders)
    
    # No cached image, fetch a new one only if creating a new presentation
    if is_new_presentation:
        return fetch_unique_image(slide_type, context, used_images, use_placeholders)
    else:
        # For cached presentations, don't fetch new images if not in cache
        return None

def fetch_unique_image(slide_type: str, context: Dict[str, Any], used_images: set, use_placeholders: bool = False):
    """Fetch a unique image that hasn't been used in the presentation yet."""
    # Don't fetch new images if we're not creating a new presentation
    if not context.get('is_new_presentation', True):
        return None
        
    # Try up to 3 times to get a unique image
    for _ in range(3):
        # Fetch a new image
        image_data = fetch_image_for_slide(slide_type, context, use_placeholders)
        if not image_data:
            return None
            
        # Check if unique
        image_data.seek(0)
        img_bytes = image_data.read()
        image_data.seek(0)
        
        import hashlib
        img_hash = hashlib.md5(img_bytes).hexdigest()
        
        if img_hash not in used_images:
            # Mark as used and cache
            used_images.add(img_hash)
            
            # Cache the image
            import streamlit as st
            if 'original_images_cache' in st.session_state:
                st.session_state.original_images_cache[f'{slide_type}_slide'] = img_bytes
            
            return image_data
    
    # If we couldn't get a unique image after 3 tries, use a generic one
    return None

def get_search_terms_for_slide(slide_type: str, context: Dict[str, Any]) -> List[str]:
    """Generate specific search terms for each slide type based on content."""
    search_terms = []
    product_name = context.get("metadata", {}).get("product_name", "") if context else ""
    
    if slide_type == "problem":
        # Get keywords from problem statement
        if "problem_slide" in context and "pain_points" in context["problem_slide"]:
            pain_points = context["problem_slide"]["pain_points"]
            for point in pain_points[:2]:  # Use first two pain points for keywords
                # Extract key nouns from the pain point
                words = [w for w in point.split() if len(w) > 4 and w.lower() not in 
                        ["their", "there", "these", "those", "about", "because", "through"]]
                search_terms.extend(words[:2])  # Add up to 2 keywords per pain point
        search_terms.extend(["problem", "challenge", "frustration"])
        
    elif slide_type == "solution":
        # Get keywords from solution description
        if "solution_slide" in context and "paragraph" in context["solution_slide"]:
            solution_text = context["solution_slide"]["paragraph"]
            # Extract key nouns from solution
            words = [w for w in solution_text.split() if len(w) > 5 and w.lower() not in 
                    ["their", "there", "these", "those", "about", "because", "through"]]
            search_terms.extend(words[:3])  # Add up to 3 keywords
        search_terms.extend([product_name, "solution", "innovation"])
        
    elif slide_type == "advantage":
        # Get keywords from advantages
        if "advantage_slide" in context and "differentiators" in context["advantage_slide"]:
            advantages = context["advantage_slide"]["differentiators"]
            for adv in advantages[:2]:  # Use first two advantages for keywords
                if isinstance(adv, str):
                    words = adv.split()
                    search_terms.extend([w for w in words if len(w) > 5][:1])
        search_terms.extend(["advantage", "unique", "competitive edge"])
        
    elif slide_type == "audience":
        # Get keywords from audience description
        if "audience_slide" in context and "paragraph" in context["audience_slide"]:
            audience_text = context["audience_slide"]["paragraph"]
            # Extract potential demographic keywords
            demographics = ["business", "enterprise", "consumer", "professional", 
                           "corporate", "small business", "tech", "healthcare", 
                           "finance", "retail", "education"]
            for demo in demographics:
                if demo in audience_text.lower():
                    search_terms.append(demo)
        search_terms.extend(["target audience", "customer", "user", "client"])
        
    elif slide_type == "market":
        # Use market-related terms
        if "market_slide" in context:
            market_desc = context["market_slide"].get("description", "")
            # Look for industry keywords
            industries = ["tech", "healthcare", "finance", "retail", "education", 
                         "manufacturing", "automotive", "energy", "telecom"]
            for ind in industries:
                if ind in market_desc.lower():
                    search_terms.append(ind)
        search_terms.extend(["market growth", "market opportunity", "business chart", "market analysis"])
    
    # Ensure we have at least some search terms
    if not search_terms:
        if product_name:
            search_terms = [product_name, slide_type]
        else:
            search_terms = [slide_type, "business"]
    
    # Clean up search terms - remove special characters and short words
    cleaned_terms = []
    for term in search_terms:
        if isinstance(term, str):
            # Remove special characters
            cleaned = ''.join(c for c in term if c.isalnum() or c.isspace())
            # Only include if substantive
            if cleaned and len(cleaned) > 3:
                cleaned_terms.append(cleaned)
    
    # Return unique search terms, with an upper limit
    unique_terms = list(set(cleaned_terms))
    return unique_terms[:5]  # Limit to 5 terms for better focus

# Improved text handling functions

def fit_text_to_box(text_frame, text, font_name=FONTS["body"], size=FONT_SIZES["body"], 
                   bold=False, color=COLORS["black"], alignment=PP_ALIGN.LEFT):
    """Format and fit text properly within a text box, adjusting size if needed."""
    text_frame.word_wrap = True
    text_frame.auto_size = False
    
    p = text_frame.paragraphs[0]
    p.alignment = alignment
    p.text = text
    
    for run in p.runs:
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)
    
    # Check if text might be too long (rough estimation)
    if len(text) > 300 and size > 16:
        # Reduce font size for very long text
        for run in p.runs:
            run.font.size = Pt(size - 2)

def create_problem_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the problem statement slide with improved text handling."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Add title with clear positioning
    title_left = Inches(0.5)  # Consistent margin
    title_top = Inches(0.5)   # Consistent position from top
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    
    # Replace any placeholder text and truncate if needed
    clean_title = content["title"].replace("[Product Name]", product_name)
    title_p.text = truncate_text_for_slide(clean_title, max_chars=60)
    title_p.alignment = PP_ALIGN.LEFT  # Left aligned for content slides
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)  # Consistent title size
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Reduce width to ensure content doesn't overflow when image is present
    content_width = Inches(6.8)  # Narrower to avoid potential overflow
    
    # Create text box for bullet points with better margins
    left = Inches(0.5)
    top = Inches(1.7)
    width = content_width
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    # Clear any existing text
    if tf.paragraphs:
        p = tf.paragraphs[0]
        p.text = ""
    else:
        p = tf.add_paragraph()
        
    # Handle different field names from slide content generators
    bullets = []
    if "pain_points" in content:
        bullets = content["pain_points"]
    elif "bullets" in content:
        bullets = content["bullets"]
    elif "differentiators" in content:
        bullets = [item["point"] if isinstance(item, dict) else str(item) for item in content["differentiators"]]
    
    # Keep max 3-4 bullets to prevent crowding and ensure spacing
    max_bullets = 4
    bullets = bullets[:max_bullets]
    
    # More aggressive truncation for bullet text
    max_bullet_chars = 75  # Reduced character limit
    max_bullet_words = 15  # Reduced word limit
    
    # Apply better spacing between bullets based on bullet count
    bullet_count = len(bullets)
    if bullet_count <= 2:
        bullet_spacing = Pt(20)  # More space for few bullets
    elif bullet_count == 3:
        bullet_spacing = Pt(15)
    else:
        bullet_spacing = Pt(12)  # Less space but still adequate for 4+ bullets
    
    for i, bullet in enumerate(bullets):
        # Replace product name placeholders
        clean_bullet = bullet.replace("[Product Name]", product_name) if product_name else bullet
        
        # Truncate bullet points with stricter limits
        bullet_text = truncate_text_for_slide(clean_bullet, max_chars=max_bullet_chars, max_words=max_bullet_words)
            
        if i == 0:
            p.text = f"• {bullet_text}"
        else:
            p = tf.add_paragraph()
            p.text = f"• {bullet_text}"
        
        # Apply improved bullet formatting
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_before = bullet_spacing
        p.space_after = bullet_spacing
        # Adjust line spacing based on bullet count
        p.line_spacing = 1.1 if bullet_count > 3 else 1.2
        
        # Apply font formatting directly to runs
        for run in p.runs:
            run.font.name = FONTS["body"]
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Add image with proper positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("problem", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        # Position image on right side - slightly adjusted for better balance
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),   # Move slightly left to avoid edge
            Inches(2.0),   # Centered vertically
            Inches(4.5),   # Consistent width
            Inches(4.0)    # Consistent height
        )
    else:
        # Add fallback icon with better positioning
        icon_left = Inches(9.5)
        icon_top = Inches(3.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_info = get_slide_icon("problem")
        icon_p.text = icon_info["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=64, alignment=PP_ALIGN.CENTER)

def create_solution_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the solution overview slide with improved text containment."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    
    # Replace any placeholder text and truncate
    clean_title = content["title"].replace("[Product Name]", product_name)
    title_p.text = truncate_text_for_slide(clean_title, max_chars=60)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Create custom text box for paragraph content - narrower width for better text fitting
    content_width = Inches(6.8)  # Reduced width to prevent overflow
    
    left = Inches(0.5)
    top = Inches(1.7)
    width = content_width
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    p = tf.paragraphs[0]
    
    # Handle different field names for content
    paragraph_text = ""
    if "paragraph" in content:
        paragraph_text = content["paragraph"]
    elif "description" in content:
        paragraph_text = content["description"]
    elif "value_proposition" in content:
        paragraph_text = content["value_proposition"]
    
    # Replace product name placeholders
    clean_paragraph = paragraph_text.replace("[Product Name]", product_name) if product_name else paragraph_text
    
    # More aggressive truncation for better text fitting
    p.text = truncate_text_for_slide(clean_paragraph, max_chars=350, max_words=70)
    
    # Format paragraph with proper spacing
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    p.line_spacing = 1.2
    
    # Apply paragraph formatting directly to runs
    for run in p.runs:
        run.font.name = FONTS["body"]
        run.font.size = Pt(18)  # Slightly smaller for better fit
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Add image with consistent positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("solution", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        # Position image on right side with consistent sizing
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),  # Adjusted position
            Inches(2.0),
            Inches(4.5),
            Inches(4.0)
        )
    else:
        # Add fallback icon
        icon_left = Inches(9.5)
        icon_top = Inches(3.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("solution")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=64, alignment=PP_ALIGN.CENTER)

def create_features_slide(prs: Presentation, content: Dict[str, Any]):
    """Create the key features slide with improved spacing."""
    slide_layout = prs.slide_layouts[6]  # Blank layout for consistency
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = truncate_text_for_slide(content["title"].replace("[Product Name]", ""), max_chars=60)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Get features list
    features = content["features"]
    max_features = 5  # Limit to 5 features for better spacing
    features = features[:max_features]
    
    # Calculate spacing based on number of features to prevent overlap
    if len(features) <= 3:
        feature_height = Inches(1.5)  # More space for few features
    else:
        feature_height = Inches(1.2)  # Less space but still adequate for more features
    
    # Create a more visually appealing list with proper spacing
    left_margin = Inches(1.0)
    top_start = Inches(1.7)
    icon_width = Inches(0.8)
    text_width = Inches(11.0)
    
    for i, feature in enumerate(features):
        # Handle different feature formats
        if isinstance(feature, dict):
            feature_text = feature.get('feature', feature.get('name', feature.get('title', str(feature))))
            icon_text = feature_text
        else:
            feature_text = str(feature)
            icon_text = feature_text
        
        # Truncate feature text appropriately
        feature_text = truncate_text_for_slide(feature_text, max_chars=100, max_words=20)

        # Calculate vertical position with dynamic spacing
        top_position = top_start + (i * feature_height)
        
        # Add icon
        icon_box = slide.shapes.add_textbox(left_margin, top_position, icon_width, feature_height)
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = match_icon_to_feature(icon_text)
        icon_p.alignment = PP_ALIGN.CENTER
        
        # Apply icon formatting
        for run in icon_p.runs:
            run.font.name = FONTS["body"]
            run.font.size = Pt(24)
            run.font.bold = True
        
        # Add feature text
        text_left = left_margin + icon_width + Inches(0.3)
        text_box = slide.shapes.add_textbox(text_left, top_position, text_width, feature_height)
        text_tf = text_box.text_frame
        text_tf.word_wrap = True
        text_tf.auto_size = False
        text_p = text_tf.paragraphs[0]
        text_p.text = feature_text
        text_p.alignment = PP_ALIGN.LEFT
        
        # Set line spacing to prevent overlap
        text_p.space_before = Pt(3)
        text_p.space_after = Pt(3)
        text_p.line_spacing = 1.1
        
        # Apply text formatting
        for run in text_p.runs:
            run.font.name = FONTS["body"]
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor.from_string(COLORS["black"])

def create_advantage_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the competitive advantage slide with improved text containment."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = content["title"].replace("[Product Name]", product_name)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Create text box for bullets - narrower for better text fitting
    content_width = Inches(6.8)  # Reduced width to prevent overflow
    
    left = Inches(0.5)
    top = Inches(1.7)
    width = content_width
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    # Clear any existing text
    if tf.paragraphs:
        p = tf.paragraphs[0]
        p.text = ""
    else:
        p = tf.add_paragraph()
        
    # Handle different field names for bulleted content
    bullets = []
    if "differentiators" in content:
        # Handle various possible dictionary structures safely
        differentiators = content["differentiators"]
        for item in differentiators:
            if isinstance(item, dict):
                # Try various possible key names for the point
                if "point" in item:
                    bullets.append(item["point"])
                elif "text" in item:
                    bullets.append(item["text"])
                elif "description" in item:
                    bullets.append(item["description"])
                elif "name" in item:
                    bullets.append(item["name"])
                elif len(item) > 0:
                    # Just take the first value if we can't find a known key
                    bullets.append(next(iter(item.values())))
                else:
                    # Empty dict, add placeholder
                    bullets.append("Competitive advantage point")
            else:
                # Handle non-dict items (strings, etc.)
                bullets.append(str(item))
    elif "bullets" in content:
        bullets = content["bullets"]
    elif "pain_points" in content:
        bullets = content["pain_points"]
    
    # If we still have no bullets, add a placeholder
    if not bullets:
        bullets = ["This product offers significant competitive advantages"]
    
    # Limit bullets for better readability
    max_bullets = 4
    bullets = bullets[:max_bullets]
    
    # More aggressive truncation for better text fitting
    max_bullet_chars = 70  # Reduced character limit
    
    for i, bullet in enumerate(bullets):
        # Replace any product name placeholders
        clean_bullet = bullet.replace("[Product Name]", product_name) if product_name and isinstance(bullet, str) else str(bullet)
        
        # Truncate with stricter limits
        bullet_text = truncate_text_for_slide(clean_bullet, max_chars=max_bullet_chars, max_words=14)
        
        if i == 0:
            p.text = f"• {bullet_text}"
        else:
            p = tf.add_paragraph()
            p.text = f"• {bullet_text}"
        
        # Apply bullet formatting
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.font.name = FONTS["body"]
        p.font.size = Pt(18)  # Slightly smaller for better fit
        p.space_before = Pt(10)
        p.space_after = Pt(10)
        p.line_spacing = 1.1  # Tighter line spacing
    
    # Add image with consistent positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("advantage", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        # Position image with adjusted dimensions
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),  # Adjusted position
            Inches(2.0),
            Inches(4.5),
            Inches(4.0)
        )
    else:
        # Add fallback icon
        icon_left = Inches(9.5)
        icon_top = Inches(3.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("advantage")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=64, alignment=PP_ALIGN.CENTER)

def create_audience_slide(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """Create the target audience slide with improved text containment."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    product_name = presentation_context.get("metadata", {}).get("product_name", "")
    
    # Add title with consistent positioning
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = content["title"].replace("[Product Name]", product_name)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Create text box for content with consistent positioning - narrower for better text fitting
    content_width = Inches(6.8)  # Reduced width to prevent overflow
    
    left = Inches(0.5)
    top = Inches(1.7)
    width = content_width
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    p = tf.paragraphs[0]
    
    # Handle different field names for text content
    paragraph_text = ""
    if "paragraph" in content:
        paragraph_text = content["paragraph"]
    elif "description" in content:
        paragraph_text = content["description"]
    elif "content" in content:
        paragraph_text = content["content"]
    
    # Replace product name placeholders
    clean_paragraph = paragraph_text.replace("[Product Name]", product_name) if product_name else paragraph_text
    
    # More aggressive truncation for better text fitting
    p.text = truncate_text_for_slide(clean_paragraph, max_chars=350, max_words=70)
    
    # Format paragraph with proper spacing
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    p.line_spacing = 1.2
    p.font.name = FONTS["body"]
    p.font.size = Pt(18)  # Slightly smaller for better fit
    
    # Add image with consistent positioning
    if used_images is None:
        used_images = set()
    
    image_data = fetch_image_with_cache("audience", presentation_context, used_images, use_placeholders=False)
    
    if image_data:
        # Position image with consistent dimensions
        pic = slide.shapes.add_picture(
            image_data,
            Inches(7.8),  # Adjusted position
            Inches(2.0),
            Inches(4.5),
            Inches(4.0)
        )
    else:
        # Add fallback icon
        icon_left = Inches(9.5)
        icon_top = Inches(3.0)
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, Inches(2.0), Inches(2.0))
        icon_tf = icon_box.text_frame
        icon_tf.auto_size = False
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = get_slide_icon("audience")["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        apply_text_formatting(icon_tf, size=64, alignment=PP_ALIGN.CENTER)

def create_cta_slide(prs: Presentation, content: Dict[str, Any]):
    """Create the call to action slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Main centered title
    title_left = Inches(0.5)
    title_top = Inches(1.5)  # Higher placement
    title_width = Inches(12.33)
    title_height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = content.get("title", "Get Started Today")
    title_p.alignment = PP_ALIGN.CENTER
    
    # Apply title formatting - larger and bolder
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(36)  # Larger for emphasis
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # CTA text - centered and prominent
    cta_left = Inches(1.0)
    cta_top = Inches(3.0)  # Positioned below title
    cta_width = Inches(11.33)
    cta_height = Inches(2.0)  # Taller for multi-line CTAs
    
    cta_box = slide.shapes.add_textbox(cta_left, cta_top, cta_width, cta_height)
    cta_tf = cta_box.text_frame
    cta_tf.word_wrap = True
    cta_tf.auto_size = False
    
    cta_p = cta_tf.paragraphs[0]
    cta_text = content.get("cta_text", content.get("call_to_action", "Contact us today to learn more!"))
    
    # Don't truncate CTAs - they're often short and important
    cta_p.text = cta_text
    cta_p.alignment = PP_ALIGN.CENTER
    
    # Apply CTA formatting - prominent and attention-grabbing
    for run in cta_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Add any supporting bullets/text if available
    if "bullets" in content and content["bullets"]:
        bullet_top = Inches(5.0)  # Position below CTA
        bullet_height = Inches(1.5)
        
        bullet_box = slide.shapes.add_textbox(cta_left, bullet_top, cta_width, bullet_height)
        bullet_tf = bullet_box.text_frame
        bullet_tf.word_wrap = True
        
        # Add up to 2 bullets only
        max_bullets = min(2, len(content["bullets"]))
        for i in range(max_bullets):
            if i == 0:
                bullet_p = bullet_tf.paragraphs[0]
            else:
                bullet_p = bullet_tf.add_paragraph()
                
            bullet_p.text = f"• {content['bullets'][i]}"
            bullet_p.alignment = PP_ALIGN.CENTER
            
            # Format bullet text
            for run in bullet_p.runs:
                run.font.name = FONTS["body"]
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Optional footer bar
    footer_left = Inches(0)
    footer_top = Inches(6.5)
    footer_width = Inches(13.33)
    footer_height = Inches(1.0)
    
    footer_box = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
    footer_box.fill.solid()
    footer_box.fill.fore_color.rgb = RGBColor.from_string(COLORS["gray"])
    footer_box.line.fill.background()

def create_market_slide_wrapper(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """
    Wrapper function for create_market_slide to ensure backward compatibility and error handling.
    """
    try:
        # Import the function from ppt_generator_additions
        from ppt_generator_additions import create_market_slide
        return create_market_slide(prs, content, presentation_context, used_images)
    except Exception as e:
        print(f"Error creating market slide: {e}")
        # Fallback to a simple slide if there's an error
        return create_fallback_slide(prs, "Market Analysis", content)

def create_roadmap_slide_wrapper(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """
    Wrapper function for create_roadmap_slide to ensure backward compatibility and error handling.
    """
    try:
        # Import the function from ppt_generator_additions
        from ppt_generator_additions import create_roadmap_slide
        return create_roadmap_slide(prs, content, presentation_context, used_images)
    except Exception as e:
        print(f"Error creating roadmap slide: {e}")
        # Fallback to a simple slide if there's an error
        return create_fallback_slide(prs, "Product Roadmap", content)

def create_team_slide_wrapper(prs: Presentation, content: Dict[str, Any], presentation_context: Dict[str, Any] = None, used_images: set = None):
    """
    Wrapper function for create_team_slide to ensure backward compatibility and error handling.
    """
    try:
        # Import the function from ppt_generator_additions
        from ppt_generator_additions import create_team_slide
        return create_team_slide(prs, content, presentation_context, used_images)
    except Exception as e:
        print(f"Error creating team slide: {e}")
        # Fallback to a simple slide if there's an error
        return create_fallback_slide(prs, "Our Team", content)

def create_fallback_slide(prs: Presentation, title: str, content: Dict[str, Any]):
    """
    Create a simple fallback slide when specialized slide generation fails.
    
    Args:
        prs: Presentation object
        title: Title to display on the slide
        content: Original content dictionary for the slide
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.33)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.auto_size = False
    title_p = title_tf.paragraphs[0]
    title_p.text = content.get("title", title)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Apply title formatting
    for run in title_p.runs:
        run.font.name = FONTS["title"]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    # Content box
    left = Inches(0.5)
    top = Inches(1.7)
    width = Inches(12.33)
    height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.auto_size = False
    
    p = tf.paragraphs[0]
    
    # Try to extract some useful content from the content dict
    text = ""
    if "description" in content:
        text += content["description"]
    elif "paragraph" in content:
        text += content["paragraph"]
    elif "content" in content:
        text += content["content"]
    
    if not text and "bullets" in content and isinstance(content["bullets"], list):
        for bullet in content["bullets"][:5]:  # Limit to 5 items
            text += f"• {bullet}\n"
    
    if not text:
        text = "Content for this slide is being prepared."
    
    p.text = text
    
    # Format text
    for run in p.runs:
        run.font.name = FONTS["body"]
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor.from_string(COLORS["black"])
    
    return slide